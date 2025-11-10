# server.py
import os, re, time, io, zipfile, uuid, json
from pathlib import Path
from typing import List, Optional, Dict

import httpx
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Pt
from dotenv import load_dotenv

# ---------- Verzeichnisse ----------
BASE_DIR = Path(__file__).resolve().parent
PUBLIC_DIR = BASE_DIR / "static"
PUBLIC_DIR.mkdir(exist_ok=True)
BUNDLES_DIR = BASE_DIR / "bundles"
BUNDLES_DIR.mkdir(exist_ok=True)

# ---------- Provider-Konfiguration ----------
load_dotenv()
OLLAMA_API_KEY    = os.getenv("OLLAMA_API_KEY", "").strip()
# OpenAI-kompatibel, z. B. https://<host>/v1
OLLAMA_CLOUD_BASE = os.getenv("OLLAMA_CLOUD_BASE", "https://ollama.com/v1").rstrip("/")

# ---------- Modell-Presets ----------
MODEL_PRESETS = {
    "deepseek-v3.1:671b-cloud": {"context_window": 131072, "ideal_max": 5000, "cap": 12000, "temperature": 0.30},
    "qwen3-coder:480b-cloud":   {"context_window": 131072, "ideal_max": 3600, "cap": 10000, "temperature": 0.25},
    "glm-4.6:cloud":            {"context_window": 131072, "ideal_max": 3200, "cap": 9000,  "temperature": 0.35},
    "gpt-oss:120b-cloud":       {"context_window": 65536,  "ideal_max": 2800, "cap": 8000,  "temperature": 0.35},
    "qwen3-vl:235b-cloud":      {"context_window": 262144, "ideal_max": 3200, "cap": 9000,  "temperature": 0.35},
    "minimax-m2:cloud":         {"context_window": 200000, "ideal_max": 2400, "cap": 8000,  "temperature": 0.35},
    "gpt-oss:20b-cloud":        {"context_window": 32768,  "ideal_max": 1400, "cap": 4000,  "temperature": 0.45},
}
MODEL_ALIASES = {
    "deepseek": "deepseek-v3.1:671b-cloud",
    "qwen3-coder": "qwen3-coder:480b-cloud",
    "glm-4.6": "glm-4.6:cloud",
    "gpt-oss:120b": "gpt-oss:120b-cloud",
    "qwen3-vl": "qwen3-vl:235b-cloud",
    "minimax-m2": "minimax-m2:cloud",
    "gpt-oss:20b": "gpt-oss:20b-cloud",
}

def resolve_model(name: str) -> str:
    n = (name or "").strip().lower()
    for key, canonical in MODEL_ALIASES.items():
        if key in n:
            return canonical
    return name if name in MODEL_PRESETS else "qwen3-coder:480b-cloud"

def choose_tokens_and_temp(model: str, requested_max: Optional[int], req_temp: Optional[float]):
    canon = resolve_model(model)
    preset = MODEL_PRESETS.get(canon, MODEL_PRESETS["qwen3-coder:480b-cloud"])
    cap = int(preset["cap"]); ideal = int(preset["ideal_max"])
    chosen_max = max(600, min(int(requested_max or ideal), cap))
    temperature = float(req_temp if req_temp is not None else preset["temperature"])
    meta = {
        "model_canonical": canon,
        "context_window": preset["context_window"],
        "ideal_max": ideal,
        "cap": cap,
        "chosen_max": chosen_max,
        "temperature": temperature,
    }
    return chosen_max, temperature, meta

# ---------- App ----------
app = FastAPI(title="AI Hub")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["GET","POST","OPTIONS"],
    allow_headers=["*"],
)
app.mount("/static", StaticFiles(directory=str(PUBLIC_DIR)), name="static")

# HTML-Seiten (Hub + Tools)
INDEX_FILE = BASE_DIR / "html/index.html"
WEBSITE_FILE = BASE_DIR / "html/website.html"
PPT_FILE = BASE_DIR / "html/ppt.html"

@app.get("/", response_class=HTMLResponse)
def root():
    if INDEX_FILE.exists():
        return INDEX_FILE.read_text(encoding="utf-8")
    return HTMLResponse("<h1>index.html fehlt</h1>", status_code=404)

@app.get("/website", response_class=HTMLResponse)
def website_page():
    if WEBSITE_FILE.exists():
        return WEBSITE_FILE.read_text(encoding="utf-8")
    return HTMLResponse("<h1>website.html fehlt</h1>", status_code=404)

@app.get("/ppt", response_class=HTMLResponse)
def ppt_page():
    if PPT_FILE.exists():
        return PPT_FILE.read_text(encoding="utf-8")
    return HTMLResponse("<h1>ppt.html fehlt</h1>", status_code=404)

@app.get("/health")
def health():
    return {"ok": True, "api_key_set": bool(OLLAMA_API_KEY), "base": OLLAMA_CLOUD_BASE}

# ---------- Schemas ----------
class GenReq(BaseModel):
    prompt: str
    model: Optional[str] = "qwen3-coder:480b-cloud"
    max_tokens: Optional[int] = None
    temperature: Optional[float] = None
    bundle_id: Optional[str] = None
    image_names: Optional[List[str]] = None

class PptReq(BaseModel):
    topic: str
    target: Optional[str] = "Allgemeines Publikum"
    slides: Optional[int] = 10
    model: Optional[str] = "qwen3-coder:480b-cloud"
    temperature: Optional[float] = 0.30

# ---------- JSON-Schemata ----------
SLIDE_JSON_SCHEMA = {
    "name": "ppt_outline",
    "schema": {
        "type": "object",
        "properties": {
            "title": {"type": "string"},
            "subtitle": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                            "minItems": 3,
                            "maxItems": 6
                        }
                    },
                    "required": ["title", "bullets"]
                },
                "minItems": 1
            },
            "closing": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "bullets": {"type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 6}
                },
                "required": ["title", "bullets"]
            }
        },
        "required": ["title", "slides", "closing"],
        "additionalProperties": False
    }
}

RESEARCH_JSON_SCHEMA = {
    "name": "ppt_research",
    "schema": {
        "type": "object",
        "properties": {
            "background": {"type": "string"},
            "key_points": {"type": "array", "items": {"type": "string"}, "minItems": 8, "maxItems": 24},
            "data_points": {"type": "array", "items": {"type": "string"}, "minItems": 6, "maxItems": 24},
            "slide_suggestions": {"type": "array", "items": {"type": "string"}},
            "sources": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {"title":{"type":"string"}, "url":{"type":"string"}},
                    "required": ["title","url"]
                }
            }
        },
        "required": ["background","key_points","data_points","slide_suggestions"],
        "additionalProperties": False
    }
}

# ---------- Helpers ----------
SAFE_NAME_RE = re.compile(r"[^A-Za-z0-9._-]+")
def safe_name(name: str) -> str:
    name = SAFE_NAME_RE.sub("", name.strip().replace(" ", "_"))
    return name.lstrip(".").replace("/", "").replace("\\", "") or f"file_{uuid.uuid4().hex[:8]}"

def strip_fences(txt: str) -> str:
    if not txt: return txt
    txt = re.sub(r"^\s*```[a-zA-Z0-9]*\s*", "", txt.strip())
    txt = re.sub(r"\s*```\s*$", "", txt)
    return txt.strip()

async def call_provider(payload: dict) -> dict:
    if not OLLAMA_API_KEY:
        raise HTTPException(status_code=500, detail="OLLAMA_API_KEY fehlt")
    url = f"{OLLAMA_CLOUD_BASE}/chat/completions"
    headers = {"Authorization": f"Bearer {OLLAMA_API_KEY}", "Content-Type": "application/json"}
    limits  = httpx.Limits(max_keepalive_connections=4, max_connections=8)
    timeout = httpx.Timeout(connect=10.0, read=240.0, write=40.0, pool=30.0)
    retriable = {408, 409, 429, 502, 503, 504}
    async with httpx.AsyncClient(http2=False, limits=limits, timeout=timeout) as client:
        backoff = 1.2
        for attempt in range(3):
            try:
                r = await client.post(url, headers=headers, json=payload)
                if r.status_code in retriable and attempt < 2:
                    time.sleep(backoff); backoff *= 1.6; continue
                if r.status_code >= 400:
                    raise HTTPException(status_code=502, detail=f"Provider {r.status_code}: {r.text[:600]}")
                return r.json()
            except httpx.RequestError as e:
                if attempt < 2:
                    time.sleep(1.5); continue
                raise HTTPException(status_code=502, detail=f"Netzwerkfehler: {e}")

def ensure_bundle(bundle_id: Optional[str]) -> str:
    bid = bundle_id or uuid.uuid4().hex[:12]
    (BUNDLES_DIR / bid / "assets").mkdir(parents=True, exist_ok=True)
    return bid

def write_html(bundle_id: str, html: str) -> Path:
    out = BUNDLES_DIR / bundle_id / "index.html"
    out.write_text(html, encoding="utf-8")
    return out

def fix_img_paths_relative(html: str, image_names: List[str]) -> str:
    for name in image_names:
        base = name.split("/")[-1]
        html = re.sub(rf'(["\'(]){re.escape(base)}([)"\'])', rf'\1assets/{base}\2', html)
    return html

def absolutize_for_preview(html: str, bundle_id: str) -> str:
    return re.sub(r'(["\'(])assets/', rf'\1/bundles/{bundle_id}/assets/', html)

def _dedupe_keep_order(seq):
    seen=set(); out=[]
    for s in seq:
        key = re.sub(r"\s+"," ",str(s)).strip().lower()
        if not key or key in seen:
            continue
        seen.add(key); out.append(str(s).strip())
    return out

def normalize_outline(outline: dict, want: int) -> dict:
    outline = outline or {}
    outline.setdefault("title", "Präsentation")
    outline.setdefault("subtitle", "")
    slides = outline.get("slides") or []
    if not isinstance(slides, list): slides = []

    cleaned=[]
    title_seen=set()
    for s in slides:
        t = str((s or {}).get("title","")).strip() or "Ohne Titel"
        tkey = t.lower()
        if tkey in title_seen:
            continue
        title_seen.add(tkey)
        bullets = _dedupe_keep_order((s or {}).get("bullets") or [])
        # generische Phrasen hart rausfiltern
        def bad(b): 
            bb=str(b).strip().lower()
            return bb in {"begriff klären","kurzes beispiel","hinweis für praxis","platzhalter","n.n."}
        bullets = [x for x in bullets if not bad(x)]
        while len(bullets) < 4:
            bullets.append("Konkreten Aspekt ergänzen")
        cleaned.append({"title": t, "bullets": bullets[:6]})
    slides = cleaned

    while len(slides) < want:
        idx = len(slides) + 1
        slides.append({"title": f"Vertiefung {idx}",
                       "bullets": ["Begriff definieren", "Messbare Kennzahl", "Beispiel mit Kontext", "Implikation"]})
    if len(slides) > want:
        slides = slides[:want]
    outline["slides"] = slides

    closing = outline.get("closing") or {}
    ctitle = str(closing.get("title","Abschluss")).strip() or "Abschluss"
    cbul = _dedupe_keep_order(closing.get("bullets") or ["Kernaussage", "Nächste Schritte"])
    while len(cbul) < 2:
        cbul.append("Nächste Schritte")
    outline["closing"] = {"title": ctitle, "bullets": cbul[:6]}
    return outline

# ---------- PPT: Bau ----------
def ppt_build(outline: dict) -> io.BytesIO:
    prs = Presentation()
    title_font_size = Pt(46)
    bullet_font_size = Pt(21)

    # Titelfolie
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = outline.get("title","Präsentation")
    if len(s.placeholders) > 1:
        s.placeholders[1].text = outline.get("subtitle","")
    try:
        s.shapes.title.text_frame.paragraphs[0].font.size = Pt(52)
        s.shapes.title.text_frame.paragraphs[0].font.bold = True
    except Exception:
        pass

    # Inhaltsfolien
    for it in outline.get("slides", []):
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = it.get("title","")
        try:
            sl.shapes.title.text_frame.paragraphs[0].font.size = title_font_size
            sl.shapes.title.text_frame.paragraphs[0].font.bold = True
        except Exception:
            pass
        body = sl.shapes.placeholders[1].text_frame
        body.clear()
        bullets = [b for b in (it.get("bullets") or []) if str(b).strip()][:6]
        if bullets:
            body.text = str(bullets[0])
            body.paragraphs[0].font.size = bullet_font_size
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = str(b)
                p.level = 0
                p.font.size = bullet_font_size

    # Abschluss
    cl = outline.get("closing")
    if cl:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = cl.get("title","Abschluss")
        try:
            sl.shapes.title.text_frame.paragraphs[0].font.size = title_font_size
            sl.shapes.title.text_frame.paragraphs[0].font.bold = True
        except Exception:
            pass
        body = sl.shapes.placeholders[1].text_frame
        body.clear()
        bullets = [b for b in (cl.get("bullets") or []) if str(b).strip()][:6]
        if bullets:
            body.text = str(bullets[0]); body.paragraphs[0].font.size = bullet_font_size
            for b in bullets[1:]:
                p = body.add_paragraph(); p.text = str(b); p.level = 0; p.font.size = bullet_font_size

    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf

# ---------- Website Generate ----------
@app.post("/generate")
async def generate(req: GenReq):
    if not req.prompt:
        raise HTTPException(status_code=400, detail="prompt fehlt")

    max_tokens, temperature, picked = choose_tokens_and_temp(req.model, req.max_tokens or None, req.temperature)
    model = picked["model_canonical"]

    bid = ensure_bundle(req.bundle_id)
    assets_dir = BUNDLES_DIR / bid / "assets"
    images_on_disk = sorted([p.name for p in assets_dir.glob("*") if p.is_file()])
    names = [n for n in (req.image_names or images_on_disk) if (assets_dir / n).exists()]

    system = (
        "Du bist ein KI-Webdesigner. Antworte NUR mit einem vollständigen, validen HTML5-Dokument "
        "inkl. eingebettetem CSS. Keine externen Skripte/Fonts.\n"
        "Wenn Bilder vorhanden sind, MUSST du sie sichtbar einbauen "
        "mit <img src=\"assets/NAME\" alt=\"…\"> in Hero, Galerie und ggf. Feature-Sektionen."
    )

    images_block = ""
    if names:
        images_block = "Verfügbare Bilder:\n" + "\n".join([f"- assets/{n}" for n in names]) + "\n"

    user = (
        f"Erstelle eine moderne One-Page basierend auf:\n\n{req.prompt}\n\n"
        f"{images_block}"
        "- Semantisches HTML, responsive Typografie und Layout, dunkles Theme erlaubt.\n"
        "- Gib ausschließlich das vollständige HTML-Dokument zurück."
    )

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user",   "content": user},
        ],
        "temperature": temperature,
        "max_tokens": max_tokens,
        "stream": False,
    }

    data = await call_provider(payload)
    content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
    html = strip_fences(content)

    if "<html" not in html.lower():
        html = (
            "<!DOCTYPE html><html lang='de'><head><meta charset='utf-8'>"
            "<meta name='viewport' content='width=device-width,initial-scale=1'>"
            "<title>Entwurf</title><style>body{font-family:Arial;padding:24px;max-width:900px;margin:0 auto}</style>"
            f"</head><body><h1>Entwurf</h1><pre>{content}</pre></body></html>"
        )

    html_saved = fix_img_paths_relative(html, names) if names else html
    html_preview = absolutize_for_preview(html_saved, bid)

    write_html(bid, html_saved)
    usage = data.get("usage", {}) if isinstance(data.get("usage", {}), dict) else {}
    return {
        "bundle_id": bid,
        "html": html_saved,
        "html_preview": html_preview,
        "meta": usage,
        "assets": names,
        "applied": picked
    }

# ---------- PPT: Recherche-Prompts ----------
def research_prompt(topic: str, target: str, slides: int) -> dict:
    sys = (
        "Du bist ein Research-Assistent. Erstelle eine faktenbasierte Zusammenfassung als JSON. "
        "Keine Floskeln, keine Platzhalter. Prägnant, prüfbar, thematisch fokussiert."
    )
    usr = (
        f"Thema: {topic}\nZielgruppe: {target}\n"
        f"Ziel: Erzeuge evidenzbasierte Inhalte für eine Präsentation mit ~{slides} Folien.\n\n"
        "JSON-Felder:\n"
        "- background: 150–300 Wörter mit Kernkontext.\n"
        "- key_points: 8–20 kurze, eigenständige Kernaussagen.\n"
        "- data_points: 6–20 knappe Fakten/Trends/Zahlen (ohne Prosa).\n"
        "- slide_suggestions: Liste mit prägnanten Folientiteln (mind. so viele wie Folien).\n"
        "- sources: optional, Liste aus Objekten {title, url}.\n"
        "Ausschließlich JSON zurückgeben."
    )
    return {"messages":[{"role":"system","content":sys},{"role":"user","content":usr}]}

def slide_prompt_from_research(topic: str, target: str, title: str, research: dict) -> dict:
    background = (research.get("background") or "")[:1800]
    key_points = research.get("key_points") or []
    data_points = research.get("data_points") or []
    kp = "\n".join(f"- {k}" for k in key_points[:20])
    dp = "\n".join(f"- {d}" for d in data_points[:20])

    sys = (
        "Du schreibst präzise Stichpunkte für Folienslides. "
        "Antworte NUR als JSON-Array von 4–6 Strings. Keine Erklärsätze, keine Meta-Hinweise, "
        "keine Floskeln wie 'Begriff klären' oder 'kurzes Beispiel'."
    )
    usr = (
        f"Thema: {topic}\nZielgruppe: {target}\nFolientitel: {title}\n\n"
        f"Kontext:\n{background}\n\n"
        f"Key Points:\n{kp}\n\nDatenpunkte:\n{dp}\n\n"
        "Gib 4–6 unterschiedlich formulierte, konkrete Bullets mit Zahlen/Beispielen, "
        "wenn sinnvoll. Vermeide Dopplungen."
    )
    return {"messages":[{"role":"system","content":sys},{"role":"user","content":usr}]}

# ---------- PPT Generate (Recherche -> Slides) ----------
@app.post("/ppt_generate")
async def ppt_generate(req: PptReq):
    if not req.topic:
        raise HTTPException(status_code=400, detail="topic fehlt")

    want = int(req.slides or 10)
    # 1) RESEARCH
    max_tokens_r, temperature_r, picked = choose_tokens_and_temp(req.model, 4000, req.temperature)
    model = picked["model_canonical"]

    payload_r = {
        "model": model,
        **research_prompt(req.topic, req.target or "Allgemeines Publikum", want),
        "temperature": max(0.1, min(0.6, temperature_r)),
        "max_tokens": max_tokens_r,
        "stream": False
    }
    # JSON erzwingen, wenn unterstützt
    payload_r["response_format"] = {"type": "json_schema", "json_schema": RESEARCH_JSON_SCHEMA}

    data_r = await call_provider(payload_r)
    txt_r = data_r.get("choices", [{}])[0].get("message", {}).get("content", "{}")
    m = re.search(r"\{.*\}\s*$", txt_r, re.S)
    raw_r = m.group(0) if m else txt_r
    try:
        research = json.loads(raw_r)
    except Exception:
        research = {"background":"", "key_points":[], "data_points":[], "slide_suggestions":[f"Folie {i+1}" for i in range(want)]}

    # Slide-Titel bestimmen
    titles = research.get("slide_suggestions") or []
    titles = [t for t in titles if str(t).strip()]
    if len(titles) < want:
        # fallback: generische, aber unterschiedliche Titel
        for i in range(len(titles), want):
            titles.append(f"Vertiefung {i+1}")

    # 2) PRO SLIDE INHALT ERZEUGEN
    slides = []
    for i in range(want):
        t = str(titles[i] if i < len(titles) else f"Vertiefung {i+1}").strip()
        # Zweiter Call pro Folie
        payload_s = {
            "model": model,
            **slide_prompt_from_research(req.topic, req.target or "Allgemeines Publikum", t, research),
            "temperature": max(0.2, min(0.7, req.temperature or 0.35)),
            "max_tokens": 600,
            "stream": False
        }
        payload_s["response_format"] = {"type": "json_object"}  # viele Anbieter akzeptieren das
        data_s = await call_provider(payload_s)
        txt_s = data_s.get("choices", [{}])[0].get("message", {}).get("content", "[]")
        m = re.search(r"\[\s*.*\s*\]\s*$", txt_s, re.S)
        raw = m.group(0) if m else txt_s
        try:
            arr = json.loads(raw)
            if isinstance(arr, dict) and "bullets" in arr:
                arr = arr["bullets"]
            if not isinstance(arr, list): arr = []
        except Exception:
            arr = []

        # Cleanup + Anti-Platzhalter
        def bad(b):
            bb=str(b).strip().lower()
            return bb in {"begriff klären","kurzes beispiel","hinweis für praxis","platzhalter","n.n."} or len(bb)<3
        bullets = [str(x).strip() for x in arr if str(x).strip() and not bad(x)]
        bullets = _dedupe_keep_order(bullets)
        if len(bullets) < 4:
            # Fallback: konstruktive Defaults
            bullets += ["Definition mit Beispiel", "Aktuelle Kennzahl/Trend", "Konkreter Einfluss/Anwendung", "Takeaway/Implikation"]
        slides.append({"title": t, "bullets": bullets[:6]})

    outline = {"title": req.topic, "subtitle": req.target or "", "slides": slides,
               "closing":{"title":"Abschluss","bullets":["Kernaussage","Nächste Schritte"]}}

    # Normalisieren auf exakt N, Duplikate entfernen
    outline = normalize_outline(outline, want)

    # 3) PPT bauen
    buf = ppt_build(outline)
    headers = {"Content-Disposition": f'attachment; filename="Slides_{safe_name(req.topic)}.pptx"'}
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers)

# ---------- ZIP ----------
@app.get("/bundle/{bundle_id}.zip")
def download_bundle(bundle_id: str):
    bundle_dir = BUNDLES_DIR / bundle_id
    if not bundle_dir.exists():
        raise HTTPException(status_code=404, detail="Bundle nicht gefunden")
    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", compression=zipfile.ZIP_DEFLATED) as z:
        for path in bundle_dir.rglob("*"):
            if path.is_file():
                z.write(path, arcname=str(path.relative_to(bundle_dir)))
    mem.seek(0)
    headers = {"Content-Disposition": f'attachment; filename="{bundle_id}.zip"'}
    return StreamingResponse(mem, media_type="application/zip", headers=headers)

# ---------- Start ----------
if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", "8080"))
    uvicorn.run("server:app", host="0.0.0.0", port=port)
